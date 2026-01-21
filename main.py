from flask import Flask, request, jsonify
from pptx import Presentation
import base64
import io
import re
from datetime import datetime

app = Flask(__name__)

def replace_placeholder_text(text_frame, replacements):
    """Replace {{placeholder}} text in a text frame with actual values"""
    if not text_frame or not text_frame.text:
        return
    
    # Get current text
    current_text = text_frame.text
    
    # Replace all {{variable}} patterns
    for key, value in replacements.items():
        pattern = '{{' + key + '}}'
        if pattern in current_text:
            current_text = current_text.replace(pattern, str(value))
    
    # Clear existing paragraphs and set new text
    text_frame.clear()
    text_frame.text = current_text

def process_shape(shape, replacements):
    """Process a single shape and replace placeholders"""
    if not hasattr(shape, 'text_frame'):
        return
    
    try:
        replace_placeholder_text(shape.text_frame, replacements)
    except Exception as e:
        print(f"Error processing shape {shape.name}: {str(e)}")

@app.route('/process-pptx', methods=['POST'])
def process_pptx():
    try:
        data = request.json
        
        # Extract binary and QBR data
        pptx_binary = data.get('pptx_binary')
        qbr_data = data.get('qbr_data', {})
        
        if not pptx_binary:
            return jsonify({'success': False, 'error': 'No PPTX binary provided'}), 400
        
        # Decode base64 PPTX
        pptx_bytes = base64.b64decode(pptx_binary)
        pptx_file = io.BytesIO(pptx_bytes)
        
        # Load presentation
        prs = Presentation(pptx_file)
        
        # Prepare replacements dictionary with all possible fields
        replacements = {
            # Config/basic info
            'period': qbr_data.get('period', ''),
            'prepared_by': qbr_data.get('preparedBy', qbr_data.get('prepared_by', '')),
            'date': datetime.now().strftime('%B %d, %Y'),
            'brand_name': qbr_data.get('clientName', qbr_data.get('brand_name', '')),
            'optional_page_refs': '',
            
            # Executive summary & insights
            'exec_summary_full': qbr_data.get('exec_summary_full', ''),
            'insight_biggest_movers': qbr_data.get('insight_biggest_movers', ''),
            'insight_below_benchmark': qbr_data.get('insight_below_benchmark', ''),
            'insight_efficiency_changes': qbr_data.get('insight_efficiency_changes', ''),
            'insight_roi_improvement': qbr_data.get('insight_roi_improvement', ''),
            
            # Analysis paragraphs
            'traffic_trends': qbr_data.get('traffic_trends', ''),
            'conversion_performance': qbr_data.get('conversion_performance', ''),
            'revenue_economics': qbr_data.get('revenue_economics', ''),
            
            # Program recommendations
            'rec_1': qbr_data.get('rec_1', ''),
            'rec_2': qbr_data.get('rec_2', ''),
            'rec_3': qbr_data.get('rec_3', ''),
            'rec_4': qbr_data.get('rec_4', ''),
            'rec_5': qbr_data.get('rec_5', ''),
            
            # Publisher performance
            'growth_drivers_paragraph': qbr_data.get('growth_drivers_paragraph', ''),
            'new_partners_paragraph': qbr_data.get('new_partners_paragraph', ''),
            'declines_paragraph': qbr_data.get('declines_paragraph', ''),
            'top_performers_paragraph': qbr_data.get('top_performers_paragraph', ''),
            'segment_insights_paragraph': qbr_data.get('segment_insights_paragraph', ''),
            
            # Publisher recommendations
            'pub_rec_1': qbr_data.get('pub_rec_1', ''),
            'pub_rec_2': qbr_data.get('pub_rec_2', ''),
            'pub_rec_3': qbr_data.get('pub_rec_3', ''),
            'pub_rec_4': qbr_data.get('pub_rec_4', ''),
            'pub_rec_5': qbr_data.get('pub_rec_5', ''),
            
            # Visibility analysis
            'brand_snapshot': qbr_data.get('brand_snapshot', ''),
            'evergreen_content': qbr_data.get('evergreen_content', ''),
            'fresh_content': qbr_data.get('fresh_content', ''),
            'discount_behavior': qbr_data.get('discount_behavior', ''),
            'category_discovery': qbr_data.get('category_discovery', ''),
            'trust_legitimacy': qbr_data.get('trust_legitimacy', ''),
            'competitors_paragraph': qbr_data.get('competitors_paragraph', ''),
            'aeo_forum_content': qbr_data.get('aeo_forum_content', ''),
            'aeo_why_forums': qbr_data.get('aeo_why_forums', ''),
            'aeo_visibility_gap': qbr_data.get('aeo_visibility_gap', ''),
            'findability_score': qbr_data.get('findability_score', '85'),
            
            # Visibility recommendations
            'vis_rec_1': qbr_data.get('vis_rec_1', ''),
            'vis_rec_2': qbr_data.get('vis_rec_2', ''),
            'vis_rec_3': qbr_data.get('vis_rec_3', ''),
            'vis_rec_4': qbr_data.get('vis_rec_4', ''),
            'vis_rec_5': qbr_data.get('vis_rec_5', ''),
            'vis_rec_6': qbr_data.get('vis_rec_6', ''),
            'vis_rec_7': qbr_data.get('vis_rec_7', ''),
            'vis_rec_8': qbr_data.get('vis_rec_8', ''),
            
            # Tables (placeholder for now)
            'yoy_summary_table': qbr_data.get('yoy_summary_table', 'See Complete Report'),
            'top_current_performers_table': qbr_data.get('top_current_performers_table', 'See Complete Report'),
            'segment_overview_table': qbr_data.get('segment_overview_table', 'See Complete Report'),
            'top_10_growth_table': qbr_data.get('top_10_growth_table', 'See Complete Report'),
            'top_10_decline_table': qbr_data.get('top_10_decline_table', 'See Complete Report'),
            'top_cited_domains_table': qbr_data.get('top_cited_domains_table', 'See Complete Report'),
            'visibility_opportunities_table': qbr_data.get('visibility_opportunities_table', 'See Complete Report'),
            
            # Metrics for snapshot slide
            'clicks_recent': qbr_data.get('Current Clicks', qbr_data.get('clicks_recent', 'See Complete Report')),
            'clicks_yoy_pct': qbr_data.get('YoY Clicks Change', qbr_data.get('clicks_yoy_pct', 'See Complete Report')),
            'sales_recent': qbr_data.get('Current Sales', qbr_data.get('sales_recent', 'See Complete Report')),
            'sales_yoy_pct': qbr_data.get('YoY Sales Change', qbr_data.get('sales_yoy_pct', 'See Complete Report')),
            'conv_rate_recent': qbr_data.get('Current Conv Rate', qbr_data.get('conv_rate_recent', 'See Complete Report')),
            'conv_rate_yoy_pct': qbr_data.get('YoY Conv Rate Change', qbr_data.get('conv_rate_yoy_pct', 'See Complete Report')),
            'order_value_recent': qbr_data.get('Current Order Value', qbr_data.get('order_value_recent', 'See Complete Report')),
            'order_value_yoy_pct': qbr_data.get('YoY Order Value Change', qbr_data.get('order_value_yoy_pct', 'See Complete Report')),
            'aov_recent': qbr_data.get('Current AOV', qbr_data.get('aov_recent', 'See Complete Report')),
            'aov_yoy_pct': qbr_data.get('YoY AOV Change', qbr_data.get('aov_yoy_pct', 'See Complete Report')),
            'pub_commission_recent': qbr_data.get('Current Commission', qbr_data.get('pub_commission_recent', 'See Complete Report')),
            'pub_commission_yoy_pct': qbr_data.get('YoY Commission Change', qbr_data.get('pub_commission_yoy_pct', 'See Complete Report')),
            'cpa_recent': qbr_data.get('Current CPA', qbr_data.get('cpa_recent', 'See Complete Report')),
            'cpa_yoy_pct': qbr_data.get('YoY CPA Change', qbr_data.get('cpa_yoy_pct', 'See Complete Report')),
            'roi_recent': qbr_data.get('Current ROI', qbr_data.get('roi_recent', 'See Complete Report')),
            'roi_yoy_pct': qbr_data.get('YoY ROI Change', qbr_data.get('roi_yoy_pct', 'See Complete Report')),
        }
        
        # Process all slides and shapes
        for slide in prs.slides:
            for shape in slide.shapes:
                process_shape(shape, replacements)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Encode to base64
        output_base64 = base64.b64encode(output.read()).decode('utf-8')
        
        return jsonify({
            'success': True,
            'binary': output_base64
        })
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"Error processing PPTX: {error_trace}")
        return jsonify({
            'success': False,
            'error': str(e),
            'trace': error_trace
        }), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy'})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000)
